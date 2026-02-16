import os
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import pytesseract
from tqdm import tqdm


# -------------------------
# FILE TYPE DETECTION
# -------------------------

def detect_file_type(file_path):
    ext = file_path.lower().split('.')[-1]

    if ext == "pdf":
        return "pdf"
    elif ext in ["jpg", "jpeg", "png"]:
        return "image"
    elif ext == "pptx":
        return "ppt"
    else:
        return "unknown"


# -------------------------
# PDF EXTRACTION
# -------------------------

import fitz
import pdfplumber

def extract_pdf(file_path):
    text = ""

    # First try PyMuPDF
    try:
        doc = fitz.open(file_path)
        for page in doc:
            text += page.get_text()
        return text
    except Exception as e:
        print(f"PyMuPDF failed for {file_path}, trying pdfplumber...")

    # Fallback to pdfplumber
    try:
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() or ""
        return text
    except Exception as e:
        print(f"Both extractors failed for {file_path}")
        return ""



# -------------------------
# PPT EXTRACTION
# -------------------------

def extract_ppt(file_path):
    text = ""
    prs = Presentation(file_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return text


# -------------------------
# IMAGE OCR EXTRACTION
# -------------------------

def extract_image(file_path):
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)
    return text


# -------------------------
# UNIFIED EXTRACTION
# -------------------------

def extract_text(file_path):
    file_type = detect_file_type(file_path)

    if file_type == "pdf":
        return extract_pdf(file_path)

    elif file_type == "ppt":
        return extract_ppt(file_path)

    elif file_type == "image":
        return extract_image(file_path)

    return ""
